#!/usr/bin/env python3
"""공학연구인턴십 성과발표회 PPT — POD-Bench (조현준 / VisAGI Lab).

Built on the author's own deck as a template so the master, theme, A4 slide
size, the blue dash header rule and the 맑은 고딕 font all carry over.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

import os
HERE = os.path.dirname(os.path.abspath(__file__))
FIG  = f"{HERE}/figs"
OUT  = f"{HERE}/POD-Bench_성과발표.pptx"
TPL  = f"{HERE}/assets/template.pptx"
RULE = f"{HERE}/assets/header_rule.png"
LOGO = f"{HERE}/assets/uos_logo.png"

# ── palette lifted from the author's deck ────────────────────────────────────
NAVY   = RGBColor(0x00, 0x40, 0x94)
BLUE   = RGBColor(0x0A, 0x4D, 0x9B)
MID    = RGBColor(0x2E, 0x75, 0xB6)
LIGHT  = RGBColor(0x9D, 0xC3, 0xE6)
PALE   = RGBColor(0xE9, 0xF1, 0xFA)
PALER  = RGBColor(0xF5, 0xF8, 0xFC)
RED    = RGBColor(0xC0, 0x00, 0x00)
PURPLE = RGBColor(0x70, 0x30, 0xA0)
INK    = RGBColor(0x27, 0x25, 0x1E)
GREY   = RGBColor(0x76, 0x71, 0x71)
LGREY  = RGBColor(0xD9, 0xD9, 0xD9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

KFONT = "맑은 고딕"

SW, SH   = 10.8333, 7.5
HDR      = (0.72, 0.44, 10.13, 0.17)      # the blue dash rule, every slide
TITLE_XY = (0.64, 0.96, 9.66, 0.52)
LM       = 0.64                            # left margin
CW       = 9.55                            # content width

prs = Presentation(TPL)
BLANK = prs.slide_layouts[6]               # 빈 화면


# ── helpers ──────────────────────────────────────────────────────────────────
def _ea(run, name=KFONT):
    """python-pptx sets only the latin typeface; Korean needs <a:ea> too."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def style(run, size=16, bold=False, color=INK, font=KFONT, italic=False):
    f = run.font
    f.name, f.size, f.bold, f.italic = font, Pt(size), bold, italic
    f.color.rgb = color
    _ea(run, font)
    return run


def add_slide(page=None):
    s = prs.slides.add_slide(BLANK)
    for ph in list(s.placeholders):                 # blank layout still carries 3
        ph._element.getparent().remove(ph._element)
    s.shapes.add_picture(RULE, Inches(HDR[0]), Inches(HDR[1]),
                         Inches(HDR[2]), Inches(HDR[3]))
    if page is not None:
        tb = s.shapes.add_textbox(Inches(SW - 1.05), Inches(SH - 0.44),
                                  Inches(0.6), Inches(0.28))
        tb.text_frame.word_wrap = False
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        style(p.add_run(), 10).text = str(page)
        p.runs[0].text = str(page); style(p.runs[0], 10, color=GREY)
    return s


def textbox(s, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tb, tf


def title(s, text, sub=None, color=NAVY):
    tb, tf = textbox(s, *TITLE_XY)
    p = tf.paragraphs[0]
    style(p.add_run(), 24, True, color).text = text
    if sub:
        r = p.add_run(); style(r, 14, False, GREY).text = "   " + sub
    return tb


def para(tf, text, size=16, bold=False, color=INK, indent=0, space_before=0,
         space_after=2, first=False, italic=False, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before, p.space_after = Pt(space_before), Pt(space_after)
    if line:
        p.line_spacing = line
    if indent:
        p.level = min(indent, 4)
    p.alignment = PP_ALIGN.LEFT
    style(p.add_run(), size, bold, color, italic=italic).text = text
    return p


def rich(tf, parts, size=16, indent=0, space_before=0, space_after=2,
         first=False, align=None, line=None):
    """parts = [(text, bold, color), ...] inside one paragraph."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before, p.space_after = Pt(space_before), Pt(space_after)
    if line:
        p.line_spacing = line
    if indent:
        p.level = min(indent, 4)
    p.alignment = align or PP_ALIGN.LEFT
    for t, b, c in parts:
        style(p.add_run(), size, b, c).text = t
    return p


def card(s, x, y, w, h, fill=PALER, line=LIGHT, lw=1.0, radius=0.10, shadow=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    if not shadow:
        sp = sh._element.spPr
        el = sp.makeelement(qn("a:effectLst"), {})
        sp.append(el)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.10)
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def chip(s, x, y, w, h, text, fill, fg=WHITE, size=12, bold=True):
    sh = card(s, x, y, w, h, fill=fill, line=None, radius=0.35)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(0.01)
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style(p.add_run(), size, bold, fg).text = text
    return sh


def arrow(s, x, y, w, h=0.18, color=MID, right=True):
    sh = s.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW if right else MSO_SHAPE.DOWN_ARROW,
        Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sp = sh._element.spPr
    sp.append(sp.makeelement(qn("a:effectLst"), {}))
    return sh


def pic(s, name, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return s.shapes.add_picture(f"{FIG}/{name}", Inches(x), Inches(y), **kw)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text



def sec_head(tf, text, first=False, space_before=0):
    """the author's signature body header: bold, blue, 16pt."""
    return para(tf, text, 15.5, True, BLUE, first=first, space_before=space_before,
                space_after=3)


def lead(s, text_parts, y=1.76, size=15):
    """one-line 'why this slide exists' sentence under the title."""
    tb, tf = textbox(s, LM, y, CW, 0.42)
    rich(tf, text_parts, size=size, first=True, line=1.2)
    return tb



# ═════════════════════════════════════════════════════════════════════════════
# 1 — 표지
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
tb, tf = textbox(s, LM, 4.10, 9.6, 0.90)
rich(tf, [("공학연구인턴십 ", True, NAVY), ("성과발표", True, NAVY)], size=40, first=True)

tb, tf = textbox(s, LM, 5.06, 9.7, 0.95)
rich(tf, [("POD-Bench", True, BLUE),
          ("  ·  AI는 ", False, INK), ("내 취향", True, INK),
          ("과 ", False, INK), ("지금 상황", True, INK),
          ("을 동시에 만족시킬 수 있는가?", False, INK)],
     size=17, first=True)
para(tf, "Personalized Outfit Decision Benchmark for Vision-Language Models",
     12.5, False, GREY, space_before=4)

tb, tf = textbox(s, LM, 6.44, 4.6, 0.62)
rich(tf, [("VisAGI Lab", True, NAVY)], size=14, first=True)
rich(tf, [("인공지능학과  2023480017  조현준", True, NAVY)], size=14)

tb, tf = textbox(s, 6.6, 6.50, 2.3, 0.30)
rich(tf, [("2026. 08. 05.", False, GREY)], size=12, first=True, align=PP_ALIGN.RIGHT)
s.shapes.add_picture(LOGO, Inches(8.25), Inches(6.74), Inches(1.95), Inches(0.32))
notes(s, "[0:00] 안녕하십니까. 인공지능학과 조현준입니다. VisAGI Lab에서 진행한 "
         "공학연구인턴십 결과를 발표하겠습니다. 주제는 AI가 사용자의 취향과 "
         "그때그때의 상황을 동시에 만족시킬 수 있는지를 재는 시험지를 만드는 일이었습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 2 — 목차
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(2)
title(s, "목차")

items = [("01", "연구실 · 인턴십 개요", "0:20"),
         ("02", "무엇을 풀려고 했나 · 어떻게 채점하나", "0:50"),
         ("03", "만드는 과정 — 다섯 단계", "2:00"),
         ("04", "결과", "1:00"),
         ("05", "기여 · 향후 계획 · 소감", "0:50")]

y = 2.10
for i, (num_, name, t) in enumerate(items):
    card(s, LM, y, CW, 0.80, fill=PALER if i % 2 == 0 else WHITE,
         line=LIGHT, radius=0.16)
    chip(s, LM + 0.20, y + 0.21, 0.46, 0.38, num_, NAVY, size=13)
    tb, tf = textbox(s, LM + 0.88, y + 0.22, 6.6, 0.40)
    rich(tf, [(name, True, BLUE)], size=16.5, first=True)
    tb, tf = textbox(s, LM + CW - 1.10, y + 0.24, 0.90, 0.34)
    rich(tf, [(t, True, MID)], size=13, first=True, align=PP_ALIGN.RIGHT)
    y += 0.90

card(s, LM, 6.50, CW, 0.52, fill=PALE, line=None, radius=0.10)
tb, tf = textbox(s, LM + 0.20, 6.60, CW - 0.40, 0.34)
rich(tf, [("이번 발표의 중심은 ", False, INK), ("③ 만드는 과정", True, NAVY),
          ("입니다 — 다섯 단계 각각이 왜 필요한지를 예시로 설명드립니다.", False, INK)],
     size=13, first=True)
notes(s, "[0:15] 발표는 다섯 부분입니다. 이 중 중심은 세 번째, 시험지를 만드는 "
         "다섯 단계입니다. 각 단계가 왜 필요한지를 예시와 함께 말씀드리겠습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 3 — 연구실 · 인턴십 개요
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(3)
title(s, "연구실 · 인턴십 개요")

lead(s, [("VisAGI Lab", True, BLUE),
         ("  ·  AI가 이미지와 언어를 함께 이해하는 능력을 ", False, INK),
         ("평가하고 개선", True, INK), ("하는 연구실", False, INK)])

stats = [("2,621", "만든 문항"), ("48", "만든 대화"),
         ("1,661", "직접 검수한 칸"), ("8종", "평가한 AI 모델")]
x = LM
for v, lab in stats:
    c = card(s, x, 2.40, 2.24, 1.02, fill=PALE, line=None, radius=0.12)
    tf2 = c.text_frame
    rich(tf2, [(v, True, NAVY)], size=30, first=True, align=PP_ALIGN.CENTER)
    rich(tf2, [(lab, False, INK)], size=11.5, align=PP_ALIGN.CENTER)
    x += 2.34

tb, tf = textbox(s, LM, 3.80, CW, 0.4)
sec_head(tf, "제가 담당한 범위", first=True)

work = [("설계", "무엇을 어떻게 잴지 정하고, 문항·대화 생성 규칙 작성"),
        ("구현", "다섯 단계 파이프라인 전체 개발"),
        ("검수", "1,661칸 검수 도구 제작 및 전수 확인"),
        ("실험", "AI 모델 8종 평가 실행과 결과 분석")]
y = 4.24
for name, desc in work:
    chip(s, LM, y, 0.86, 0.34, name, MID, size=12)
    tb, tf = textbox(s, LM + 1.04, y + 0.05, CW - 1.10, 0.30)
    rich(tf, [(desc, False, INK)], size=13.5, first=True)
    y += 0.50

card(s, LM, 6.30, CW, 0.62, fill=PALER, line=LIGHT, radius=0.10)
tb, tf = textbox(s, LM + 0.20, 6.42, CW - 0.40, 0.40)
rich(tf, [("사용 기술  ", True, GREY),
          ("PyTorch · vLLM · FAISS · FashionSigLIP · SAM3 · Qwen3-VL · GPT-5", False, GREY)],
     size=12, first=True)
notes(s, "[0:20] VisAGI Lab은 AI가 이미지와 언어를 함께 이해하는 능력을 평가하는 "
         "연구실입니다. 저는 8주 동안 설계부터 파이프라인 구현, 검수, 그리고 "
         "AI 모델 8종 평가까지 전 과정을 맡았습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 4 — 무엇을 풀려고 했나
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(4)
title(s, "무엇을 풀려고 했나", sub="취향과 상황은 서로 다른 문제다")

lead(s, [("AI에게 옷을 추천해 달라고 하면, ", False, INK),
         ("내 취향", True, RED), ("과 ", False, INK),
         ("지금 상황", True, MID), ("을 ", False, INK),
         ("동시에", True, PURPLE), (" 맞출 수 있을까?", False, INK)], size=17)

card(s, LM, 2.26, 4.55, 1.24, fill=PALER, line=LIGHT, radius=0.10)
tb, tf = textbox(s, LM + 0.20, 2.40, 4.18, 1.00)
rich(tf, [("이 사람은", True, NAVY), ("   파란색을 좋아하고 주황색을 싫어합니다", False, INK)],
     size=13, first=True, line=1.2)
rich(tf, [("지금 상황은", True, NAVY), ("   한겨울 야외 시장, 하루 종일", False, INK)],
     size=13, space_before=8, line=1.2)

card(s, 5.55, 2.26, 4.64, 1.24, fill=WHITE, line=RED, radius=0.10)
tb, tf = textbox(s, 5.73, 2.40, 4.28, 1.00)
rich(tf, [("한쪽만 맞추면 이렇게 됩니다", True, RED)], size=13, first=True)
rich(tf, [("취향만", True, INK), ("  →  파란 원피스 ", False, INK), ("(한겨울에 춥다)", False, RED)],
     size=12.5, space_before=5, line=1.2)
rich(tf, [("상황만", True, INK), ("  →  주황 후드 ", False, INK), ("(싫어하는 색)", False, RED)],
     size=12.5, space_before=3, line=1.2)

arrow(s, 5.30, 3.58, 0.22, 0.22, color=LIGHT, right=False)

FIG_X, FIG_W = 1.42, 8.00
opt_labels = [("A", "둘 다 만족", PURPLE), ("B", "상황만 맞음", RED),
              ("C", "취향만 맞음", BLUE), ("D", "둘 다 아님", GREY)]
for i, (letter, lab, col) in enumerate(opt_labels):
    cx = FIG_X + FIG_W * (i + 0.5) / 4
    tb, tf = textbox(s, cx - 1.0, 3.92, 2.0, 0.32)
    rich(tf, [(letter, True, col)], size=19, first=True, align=PP_ALIGN.CENTER)

pic(s, "fig_abcd.png", FIG_X, 4.28, w=FIG_W)

for i, (letter, lab, col) in enumerate(opt_labels):
    cx = FIG_X + FIG_W * (i + 0.5) / 4
    tb, tf = textbox(s, cx - 1.0, 6.42, 2.0, 0.34)
    rich(tf, [(lab, True, col)], size=13.5, first=True, align=PP_ALIGN.CENTER)

tb, tf = textbox(s, LM, 6.94, CW, 0.30)
rich(tf, [("A만이 정답", True, NAVY),
          ("입니다 — 나머지 셋은 각각 어디서 틀렸는지를 알려 줍니다.", False, GREY)],
     size=12.5, first=True, align=PP_ALIGN.CENTER)
notes(s, "[0:40] 예를 들겠습니다. 파란색을 좋아하고 주황색을 싫어하는 사람이 "
         "한겨울 야외 시장에 하루 종일 있어야 합니다. 취향만 맞추면 파란 원피스인데 "
         "춥고, 상황만 맞추면 주황 후드인데 싫어하는 색입니다. "
         "기존 평가는 둘 다 '개인화된 답'으로 봅니다. 그래서 네 선지를 만들고 "
         "둘 다 만족하는 A만 정답으로 뒀습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 5 — 어떻게 채점하나
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(5)
title(s, "어떻게 채점하나", sub="답 하나로 두 가지를 따로 잰다")

lead(s, [("“틀렸다”가 아니라 ", False, INK), ("어느 쪽에서 틀렸는지", True, BLUE),
         ("까지 나와야 고칠 수 있습니다.", False, INK)])

opts = [("A", "좋아하는 색\n+ 따뜻한 옷", PURPLE), ("B", "싫어하는 색\n+ 따뜻한 옷", RED),
        ("C", "좋아하는 색\n+ 추운 옷", BLUE),   ("D", "싫어하는 색\n+ 추운 옷", GREY)]
scores = [("둘 다 맞춤", "진짜 정답", [1, 0, 0, 0], NAVY),
          ("상황 점수", "추위는 피했나", [1, 1, 0, 0], MID),
          ("취향 점수", "좋아하는 색인가", [1, 0, 1, 0], RGBColor(0x5B, 0x9B, 0xD5))]

TX, TY, TW = 0.90, 2.32, 9.05
tbl = s.shapes.add_table(4, 5, Inches(TX), Inches(TY), Inches(TW), Inches(2.60)).table
tbl.columns[0].width = Inches(2.65)
for c in range(1, 5):
    tbl.columns[c].width = Inches(1.60)
tbl.rows[0].height = Inches(0.88)
for r in range(1, 4):
    tbl.rows[r].height = Inches(0.57)

for r in range(4):
    for c in range(5):
        cell = tbl.cell(r, c)
        cell.margin_left = cell.margin_right = Inches(0.06)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        tf2 = cell.text_frame
        p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        if r == 0 and c == 0:
            pass
        elif r == 0:
            letter, meaning, col = opts[c - 1]
            cell.fill.fore_color.rgb = PALE if c == 1 else WHITE
            style(p.add_run(), 20, True, col).text = letter
            rich(tf2, [(meaning, False, INK)], size=10, align=PP_ALIGN.CENTER, line=1.1)
        elif c == 0:
            name, why, _h, col = scores[r - 1]
            cell.fill.fore_color.rgb = PALER
            p.alignment = PP_ALIGN.LEFT
            style(p.add_run(), 14, True, col).text = name
            rich(tf2, [(why, False, GREY)], size=10.5, align=PP_ALIGN.LEFT)
        else:
            if scores[r - 1][2][c - 1]:
                cell.fill.fore_color.rgb = PALE if r == 1 else WHITE
                style(p.add_run(), 20, True, scores[r - 1][3]).text = "●"
            else:
                style(p.add_run(), 14, False, LGREY).text = "·"

tb, tf = textbox(s, LM, 5.14, CW, 0.40)
sec_head(tf, "이렇게 하면 답 하나에서 두 가지가 따로 보입니다", first=True)

reads = [("A를 골랐다", "취향도 상황도 맞춤", NAVY),
         ("B를 골랐다", "상황은 알지만 취향을 무시", RED),
         ("C를 골랐다", "취향은 알지만 상황을 무시", MID)]
x = LM
for name, desc, col in reads:
    c = card(s, x, 5.58, 3.09, 0.74, fill=WHITE, line=LIGHT, radius=0.10)
    tf2 = c.text_frame
    rich(tf2, [(name, True, col)], size=13, first=True, align=PP_ALIGN.CENTER)
    rich(tf2, [(desc, False, INK)], size=11.5, align=PP_ALIGN.CENTER)
    x += 3.23

card(s, LM, 6.50, CW, 0.54, fill=PALER, line=LIGHT, radius=0.10)
tb, tf = textbox(s, LM + 0.20, 6.60, CW - 0.40, 0.36)
rich(tf, [("찍어서 맞을 확률이 정확히 25%", True, NAVY),
          ("가 되도록 선지를 배치했습니다 — 운으로 오른 점수는 걸러집니다.", False, INK)],
     size=12.5, first=True)
notes(s, "[1:05] 채점 방식이 이 연구의 핵심입니다. A를 고르면 둘 다 맞춘 것, "
         "B는 상황은 아는데 취향을 무시한 것, C는 그 반대입니다. "
         "그래서 '틀렸다'가 아니라 '어느 쪽에서 틀렸다'까지 나옵니다. "
         "찍어서 맞을 확률은 정확히 25%로 맞춰 뒀습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 6 — 전체 파이프라인 다섯 단계
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(6)
title(s, "만드는 과정 — 다섯 단계", sub="어느 하나도 건너뛸 수 없습니다")

lead(s, [("시험지 하나를 만들려면 ", False, INK),
         ("문항 · 취향 대화 · 사진 · 사람 검수", True, BLUE),
         ("가 모두 필요합니다.", False, INK)])

PIPE = [("①", "Construction", "문항 만들기", "24명 × 59개 상황\n→ 문항 2,621개", NAVY),
        ("②", "Dialog", "취향을 대화로", "취향 18개를 대화 속\n사진으로만 전달", RED),
        ("③", "Image collection", "사진 모으기", "색·옷·무늬가 맞는\n실제 상품 사진 찾기", MID),
        ("④", "Annotation", "사람이 확인", "1,661칸을 눈으로\n직접 확인", MID),
        ("⑤", "Eval", "모델 시험", "AI 8종에게 풀리고\n세 점수로 채점", NAVY)]

PW, GAP = 1.79, 0.16
x = LM
for i, (n_, en, ko, desc, col) in enumerate(PIPE):
    c = card(s, x, 2.36, PW, 2.10, fill=WHITE, line=col, radius=0.08)
    chip(s, x + PW/2 - 0.19, 2.50, 0.38, 0.34, n_, col, size=13)
    tb, tf = textbox(s, x + 0.10, 2.94, PW - 0.20, 0.52)
    rich(tf, [(ko, True, col)], size=13.5, first=True, align=PP_ALIGN.CENTER)
    rich(tf, [(en, False, GREY)], size=9.5, align=PP_ALIGN.CENTER)
    tb, tf = textbox(s, x + 0.10, 3.60, PW - 0.20, 0.74)
    rich(tf, [(desc, False, INK)], size=10.5, first=True, align=PP_ALIGN.CENTER, line=1.25)
    if i < 4:
        arrow(s, x + PW + 0.02, 3.32, 0.12, 0.16)
    x += PW + GAP

# the loop back: what the human approved becomes an input to stage 1
LOOP_Y = 4.72
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(LM + 0.30), Inches(LOOP_Y),
                        Inches(CW - 0.60), Inches(0.40))
sh.fill.background(); sh.line.color.rgb = RED; sh.line.width = Pt(1.25)
sh._element.spPr.append(sh._element.spPr.makeelement(qn("a:effectLst"), {}))
sh.text_frame.text = ""
ua = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(LM + 0.24), Inches(4.46),
                        Inches(0.18), Inches(0.22))
ua.fill.solid(); ua.fill.fore_color.rgb = RED; ua.line.fill.background()
ua._element.spPr.append(ua._element.spPr.makeelement(qn("a:effectLst"), {}))
tb, tf = textbox(s, LM + 0.60, 4.82, CW - 1.20, 0.28)
rich(tf, [("④에서 사람이 확인한 목록이 ", False, INK),
          ("①의 입력으로 되돌아갑니다", True, RED),
          ("  —  ‘사진이 있는 조합’만 문항이 되도록", False, INK)],
     size=12, first=True, align=PP_ALIGN.CENTER)

tb, tf = textbox(s, LM, 5.42, CW, 0.40)
sec_head(tf, "각 단계가 없으면 무엇이 깨지나", first=True)

breaks = [("① 없으면", "제 취향이 곧 정답이 되어 버립니다"),
          ("② 없으면", "취향을 목록으로만 줘서, 실제 대화 상황과 멀어집니다"),
          ("③ 없으면", "선지가 글자뿐이라 ‘보고 고르는’ 능력을 못 잽니다"),
          ("④ 없으면", "사진이 실제로 맞는지 아무도 보증하지 않습니다"),
          ("⑤ 없으면", "만들어만 놓고 무엇을 알아냈는지 말할 수 없습니다")]
y = 5.84
for lab, why in breaks:
    tb, tf = textbox(s, LM, y, CW, 0.26)
    rich(tf, [(lab, True, MID), ("   " + why, False, INK)], size=11.5, first=True)
    y += 0.26
notes(s, "[1:30] 시험지 하나를 만드는 데 다섯 단계가 들어갑니다. "
         "문항을 만들고, 취향을 대화로 전달하고, 사진을 모으고, 사람이 확인하고, "
         "모델을 시험합니다. 특히 빨간 화살표를 봐 주십시오. "
         "네 번째 단계에서 사람이 확인한 목록이 첫 단계의 입력으로 되돌아갑니다. "
         "이게 뒤에서 말씀드릴 문제 해결의 핵심이었습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 7 — ① Construction
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(7)
title(s, "①  Construction — 문항 만들기", sub="사람 → 상황 → 문항")

lead(s, [("문항을 손으로 쓰면 ", False, INK), ("제 취향이 정답이 되어 버립니다", True, RED),
         (". 그래서 규칙으로 만들었습니다.", False, INK)])

steps = [("사람 만들기", "좋아하는·싫어하는\n색·무늬·옷을 규칙으로 배정",
          "파랑 · 베이지 좋아함\n주황 · 빨강 싫어함", "24명"),
         ("상황 만들기", "누구나 같은 답을 낼\n상황만 채택",
          "한겨울 야외\n→ 후드 O, 원피스 X", "59개"),
         ("문항 만들기", "사람 × 상황을 곱해\n4지선다 하나로",
          "파란 후드 / 주황 후드\n파란 원피스 / 주황 원피스", "2,621개")]
x = LM
for i, (name, what, ex, count) in enumerate(steps):
    card(s, x, 2.34, 3.03, 3.16, fill=WHITE, line=MID, radius=0.08)
    tb, tf = textbox(s, x + 0.18, 2.48, 2.67, 0.34)
    rich(tf, [(name, True, NAVY)], size=15, first=True)
    tb, tf = textbox(s, x + 0.18, 2.88, 2.67, 0.66)
    rich(tf, [(what, False, INK)], size=12, first=True, line=1.25)
    card(s, x + 0.18, 3.66, 2.67, 1.00, fill=PALE, line=None, radius=0.10)
    tb, tf = textbox(s, x + 0.30, 3.78, 2.43, 0.78)
    rich(tf, [("예시", True, MID)], size=10, first=True)
    rich(tf, [(ex, False, INK)], size=11, space_before=2, line=1.2)
    tb, tf = textbox(s, x + 0.18, 4.86, 2.67, 0.42)
    rich(tf, [(count, True, NAVY)], size=19, first=True, align=PP_ALIGN.CENTER)
    if i < 2:
        arrow(s, x + 3.09, 3.72, 0.12, 0.16)
    x += 3.23

card(s, LM, 5.72, CW, 0.66, fill=PALER, line=LIGHT, radius=0.10)
tb, tf = textbox(s, LM + 0.20, 5.84, CW - 0.40, 0.44)
rich(tf, [("세 단계 모두 사람 손이 들어가지 않습니다", True, NAVY),
          (" — 같은 씨앗을 넣으면 언제나 똑같은 데이터가 나옵니다.", False, INK)],
     size=13, first=True)

tb, tf = textbox(s, LM, 6.56, CW, 0.44)
rich(tf, [("왜 중요한가  ", True, GREY),
          ("남이 제 결과를 다시 만들어 검증할 수 있고, 제가 데이터를 유리하게 "
           "고르지 않았다는 증거가 됩니다.", False, GREY)], size=11.5, first=True, line=1.2)
notes(s, "[1:50] 첫 단계입니다. 사람 24명의 취향을 규칙으로 배정하고, 누구나 같은 "
         "답을 낼 상황 59개를 만든 뒤, 둘을 곱해 2,621개 문항을 얻습니다. "
         "사람 손이 안 들어가서 같은 씨앗이면 항상 같은 데이터가 나옵니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 8 — ② Dialog : 왜 대화인가
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(8)
title(s, "②  Dialog — 취향을 대화로 전달", sub="이번에 새로 추가한 단계")

lead(s, [("실제 사용자는 취향을 ", False, INK),
         ("목록으로 말하지 않습니다", True, RED),
         (" — 대화 중에 사진으로 흘립니다.", False, INK)])

card(s, LM, 2.24, 4.55, 1.22, fill=WHITE, line=LGREY, radius=0.10)
tb, tf = textbox(s, LM + 0.20, 2.38, 4.18, 0.98)
rich(tf, [("지금까지 (목록·줄글)", True, GREY)], size=13, first=True)
rich(tf, [("“좋아하는 색: 파랑, 베이지 · 싫어하는 색: 주황 …”", False, INK)],
     size=11.5, space_before=5, line=1.2)
rich(tf, [("한 번에 다 알려주는 형태", False, GREY)], size=11, space_before=4)

card(s, 5.55, 2.24, 4.64, 1.22, fill=PALE, line=NAVY, radius=0.10)
tb, tf = textbox(s, 5.73, 2.38, 4.28, 0.98)
rich(tf, [("새로 만든 것 (대화)", True, NAVY)], size=13, first=True)
rich(tf, [("“이 둘은 하나만 다른데, ", False, INK),
          ("첫 번째 옷의 커트", True, NAVY), ("가 내가 절대 안 입는 커트야”", False, INK)],
     size=11.5, space_before=5, line=1.2)
rich(tf, [("+ 사진 2장  →  값은 사진에만 있음", True, RED)], size=11, space_before=4)

tb, tf = textbox(s, LM, 3.62, CW, 0.40)
sec_head(tf, "대화가 취향을 전달하는 네 가지 방식", first=True)

MODES = [("single", "1장", "그 사진의 값을 읽으면 됨", MID),
         ("differ", "2장", "한 축만 다름 — 어느 쪽인지 텍스트가 지정", MID),
         ("share2", "2장", "두 장이 공유하는 값 하나를 찾아야 함", NAVY),
         ("share3", "3장", "세 장의 교집합을 계산해야 함", NAVY)]
y = 4.06
for name, nimg, what, col in MODES:
    card(s, LM, y, CW, 0.50, fill=PALER if col == MID else PALE, line=None, radius=0.10)
    tb, tf = textbox(s, LM + 0.22, y + 0.13, 1.5, 0.28)
    rich(tf, [(name, True, col)], size=13, first=True)
    tb, tf = textbox(s, LM + 1.72, y + 0.14, 0.9, 0.26)
    rich(tf, [("사진 " + nimg, True, GREY)], size=11.5, first=True)
    tb, tf = textbox(s, LM + 2.86, y + 0.14, CW - 3.10, 0.26)
    rich(tf, [(what, False, INK)], size=12, first=True)
    y += 0.56

card(s, LM, 6.36, CW, 0.72, fill=WHITE, line=RED, radius=0.10)
tb, tf = textbox(s, LM + 0.20, 6.48, CW - 0.40, 0.50)
rich(tf, [("절대 규칙  ", True, RED),
          ("대화에는 상황 이야기를 한 마디도 넣지 않습니다.", True, INK)],
     size=13, first=True)
rich(tf, [("대화 속 상대가 “결혼식엔 흰색은 피하세요”라고 말하는 순간, "
           "재려던 능력이 대화로 새어 나갑니다.", False, GREY)],
     size=11, space_before=3)
notes(s, "[2:15] 이번에 새로 추가한 단계입니다. 지금까지는 취향을 목록으로 "
         "한 번에 줬는데, 실제 사용자는 그렇게 말하지 않습니다. "
         "그래서 취향 18개를 전부 대화 속에서, 그것도 값은 오직 사진에만 담아 "
         "전달하도록 만들었습니다. 네 가지 방식이 있고 뒤로 갈수록 어렵습니다. "
         "가장 중요한 규칙은 대화에 상황 이야기를 절대 넣지 않는 것입니다. "
         "넣는 순간 재려던 능력이 새어 나가기 때문입니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 9 — ② Dialog : 실제 대화와 검증
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(9)
title(s, "②  Dialog — 실제로 만들어진 대화", sub="사진판과 글자판은 쌍둥이")

# chat bubbles — the real turns from dlg__U001__image.json
tb, tf = textbox(s, LM, 1.74, CW, 0.30)
rich(tf, [("사진판  ", True, NAVY), ("dlg__U001__image.json  ·  실제 2번 턴", False, GREY)],
     size=12, first=True)

card(s, LM, 2.06, 6.60, 0.62, fill=PALE, line=None, radius=0.22)
tb, tf = textbox(s, LM + 0.22, 2.20, 6.20, 0.36)
rich(tf, [("“이 둘은 같은 물건인데 하나만 달라. ", False, INK),
          ("첫 번째 옷의 커트", True, NAVY),
          ("는 내가 절대 안 입는 커트야.”", False, INK)], size=12.5, first=True)

pic(s, "fig_dlg_example.png", 1.30, 2.84, w=2.50)
tb, tf = textbox(s, 0.90, 4.08, 3.30, 0.26)
rich(tf, [("사진 2장만 제시  ·  옷 이름은 어디에도 없음", False, GREY)],
     size=10.5, first=True, align=PP_ALIGN.CENTER)

card(s, 4.86, 2.80, 5.33, 1.42, fill=WHITE, line=MID, radius=0.10)
tb, tf = textbox(s, 5.04, 2.92, 4.97, 1.18)
rich(tf, [("모델이 해야 하는 일", True, MID)], size=12.5, first=True)
rich(tf, [("① 두 사진을 보고 무엇이 다른지 찾는다", False, INK)], size=11.5, space_before=4)
rich(tf, [("② ‘커트’가 옷 종류를 가리킨다는 걸 안다", False, INK)], size=11.5, space_before=2)
rich(tf, [("③ 첫 번째 사진의 옷 종류를 ", False, INK),
          ("싫어함", True, RED), ("으로 기록한다", False, INK)], size=11.5, space_before=2)

tb, tf = textbox(s, LM, 4.42, CW, 0.30)
rich(tf, [("글자판  ", True, NAVY),
          ("같은 대화, 같은 순서 — 사진 자리에만 이름이 들어갑니다", False, GREY)],
     size=12, first=True)
card(s, LM, 4.74, CW, 0.58, fill=PALER, line=LGREY, radius=0.22)
tb, tf = textbox(s, LM + 0.22, 4.86, CW - 0.44, 0.36)
rich(tf, [("“…첫 번째 옷의 커트는 내가 절대 안 입는 커트야.  ", False, INK),
          ("{navy solid pea coat}   {navy solid windbreaker}", True, MID)],
     size=12, first=True)

tb, tf = textbox(s, LM, 5.48, CW, 0.30)
rich(tf, [("두 판의 차이가 딱 하나 ", False, INK), ("— 값이 사진이냐 글자냐", True, NAVY),
          ("  →  성능 차이가 곧 ‘보는 능력’의 몫", False, INK)], size=12.5, first=True)

vals = [("48개", "만든 대화\n(사진판 24 + 글자판 24)"),
        ("432개", "전달한 취향\n24명 × 18개, 빠짐없이"),
        ("48/48", "검증 통과\n규칙 5종 전수 확인"),
        ("6/6", "일부러 망가뜨린 것을\n검증기가 모두 검출")]
x = LM
for v, lab in vals:
    c = card(s, x, 5.92, 2.24, 1.06, fill=PALE if "6/6" in v or "48/48" in v else PALER,
             line=None, radius=0.12)
    tf2 = c.text_frame
    rich(tf2, [(v, True, NAVY)], size=20, first=True, align=PP_ALIGN.CENTER)
    rich(tf2, [(lab, False, INK)], size=9.5, align=PP_ALIGN.CENTER, line=1.15)
    x += 2.34

tb, tf = textbox(s, LM, 7.06, CW, 0.28)
rich(tf, [("검증기는 생성기와 별개 프로그램입니다 — 자기가 만든 걸 자기 기준으로 "
           "통과시키면 검증이 아니기 때문입니다.", False, GREY)], size=11, first=True)
notes(s, "[2:40] 실제로 만들어진 대화입니다. 사진 두 장을 주고 '첫 번째 옷의 커트는 "
         "내가 안 입는 커트'라고만 말합니다. 옷 이름은 어디에도 없습니다. "
         "모델은 두 사진의 차이를 찾고, 커트가 옷 종류를 뜻한다는 걸 알고, "
         "첫 번째 옷을 싫어함으로 기록해야 합니다. "
         "글자판은 완전히 같은 대화인데 사진 자리에만 이름이 들어갑니다. "
         "두 판의 차이가 곧 보는 능력의 몫입니다. "
         "대화 48개 전부 검증을 통과했고, 일부러 망가뜨린 여섯 가지를 검증기가 "
         "모두 잡아냈습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 10 — ③ Image collection
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(10)
title(s, "③  Image collection — 사진 모으기", sub="문항 하나에 실제 상품 사진 4장")

lead(s, [("검색만으로는 ", False, INK), ("원하는 사진이 나오지 않습니다", True, RED),
         (" — 색·옷·무늬가 모두 맞아야 합니다.", False, INK)])

FUN = [("검색으로 모은 사진",      "166,100", 6.20, MID,   ""),
       ("옷 부분을 찾아낸 것",     "127,859", 5.50, MID,   "나머지는 옷이 안 보임"),
       ("색·옷·무늬가 맞는 것",    "38,829",  3.40, LIGHT, "셋 다 통과한 것만"),
       ("후보로 저장",            "19,012",  2.70, LIGHT, "칸마다 상위 10장"),
       ("④ 사람 확인으로 넘김",    "1,661",   1.55, NAVY,  "확인해야 할 칸")]
CX, FY, BH, GAP = 5.62, 2.40, 0.44, 0.14
for i, (name, val, w_, col, note) in enumerate(FUN):
    x0 = CX - w_ / 2
    c = card(s, x0, FY, w_, BH, fill=col, line=None, radius=0.10)
    rich(c.text_frame, [(val, True, WHITE)], size=14, first=True, align=PP_ALIGN.CENTER)
    tb, tf = textbox(s, LM, FY + 0.10, CX - w_ / 2 - LM - 0.16, 0.28)
    rich(tf, [(name, i in (0, 4), NAVY if i == 4 else INK)], size=12,
         first=True, align=PP_ALIGN.RIGHT)
    if note:
        tb, tf = textbox(s, CX + w_ / 2 + 0.16, FY + 0.11, 2.6, 0.26)
        rich(tf, [(note, False, GREY)], size=10.5, first=True)
    if i < len(FUN) - 1:
        arrow(s, CX - 0.05, FY + BH + 0.01, 0.10, 0.12,
              color=RGBColor(0xB4, 0xB4, 0xB4), right=False)
    FY += BH + GAP

tb, tf = textbox(s, LM, 5.56, CW, 0.40)
sec_head(tf, "축마다 확인 방법을 다르게 한 이유", first=True)

axes_info = [("옷 종류", "AI에게 물어봄", "옷을 오려내도 원피스와\n민소매를 구분 못 함", NAVY),
             ("무늬", "큰 칸으로 나눠 확인", "무늬는 넓게 봐야\n보이기 때문", MID),
             ("색", "작은 칸으로 나눠 확인", "남색과 파랑이 섞이는 걸\n잡아내려고", MID)]
x = LM
for name, method, why, col in axes_info:
    c = card(s, x, 6.00, 3.09, 1.06, fill=WHITE, line=LIGHT, radius=0.10)
    tf2 = c.text_frame
    rich(tf2, [(name + "  ", True, col), (method, True, INK)], size=12,
         first=True, align=PP_ALIGN.CENTER)
    rich(tf2, [(why, False, GREY)], size=10, align=PP_ALIGN.CENTER, line=1.15)
    x += 3.23
notes(s, "[3:05] 문항 하나에 사진 네 장이 필요한데, 색과 옷과 무늬가 정확히 "
         "맞아야 합니다. 16만 장에서 시작해 옷 부분만 오려내고, 세 축을 각각 "
         "다른 방법으로 확인해서 1,661칸까지 좁힙니다. "
         "축마다 방법을 다르게 한 이유가 아래에 있습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 11 — 문제해결 ① 국소 패턴
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(11)
title(s, "③에서 부딪힌 문제 — 검색이 엉뚱한 옷을 준다")

lead(s, [("“체크무늬 베이지 셔츠”를 찾으면 ", False, INK),
         ("주머니에만 체크가 있는 옷", True, RED), ("이 통과했습니다.", False, INK)])

pic(s, "fig_patch.png", 1.92, 2.14, w=7.00)

tb, tf = textbox(s, LM, 4.70, CW, 0.4)
sec_head(tf, "세 가지를 시도했고, 셋 다 실패했습니다", first=True)

tries = [("검색어를 더 자세히", "옷 정확도만 조금 올라감"),
         ("재정렬 모델 추가", "효과 없음"),
         ("AI에게 판정 맡기기", "똑같은 한계")]
x = LM
for name, res in tries:
    c = card(s, x, 5.06, 3.09, 0.68, fill=WHITE, line=LGREY, radius=0.10)
    tf2 = c.text_frame
    rich(tf2, [(name, True, GREY)], size=12.5, first=True, align=PP_ALIGN.CENTER)
    rich(tf2, [("✕  " + res, True, RED)], size=11.5, align=PP_ALIGN.CENTER)
    x += 3.23

card(s, LM, 5.88, CW, 1.08, fill=PALE, line=NAVY, radius=0.08)
tb, tf = textbox(s, LM + 0.20, 6.00, CW - 0.40, 0.86)
rich(tf, [("원인  ", True, NAVY), ("세 방법 모두 사진 한 장을 ", False, INK),
          ("숫자 하나로 요약", True, NAVY),
          ("해서 비교합니다. 무늬가 어디에 있는지가 평균에 묻힙니다.", False, INK)],
     size=12.5, first=True, line=1.2)
rich(tf, [("해결  ", True, NAVY),
          ("옷을 잘라 격자로 나눈 뒤 칸마다 따로 확인 — 모델을 새로 학습시키지 않고 "
           "보는 단위만 바꿨습니다.", False, INK)], size=12.5, space_before=5, line=1.2)
notes(s, "[3:25] 여기서 부딪힌 문제입니다. 체크무늬 셔츠를 찾으면 주머니에만 "
         "체크가 있는 옷이 통과했습니다. 세 가지를 시도했지만 다 실패했고, "
         "원인은 셋 다 사진 한 장을 숫자 하나로 요약해서 비교한다는 점이었습니다. "
         "그래서 격자로 나눠 칸마다 확인하도록 바꿔 해결했습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 12 — ④ Annotation + 되먹임
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(12)
title(s, "④  Annotation — 사람이 직접 확인", sub="그리고 그 결과가 ①로 되돌아간다")

lead(s, [("기계가 고른 것을 ", False, INK), ("사람이 한 칸씩 눈으로 확인", True, BLUE),
         ("했습니다. 1,661칸 전부.", False, INK)])

nums = [("1,661", "확인한 칸", INK), ("1,223", "쓸 수 있음 (74%)", NAVY),
        ("438", "제외", RED)]
x = LM
for v, lab, col in nums:
    c = card(s, x, 2.24, 3.09, 0.86, fill=PALE if col == NAVY else PALER,
             line=None, radius=0.12)
    tf2 = c.text_frame
    rich(tf2, [(v, True, col)], size=24, first=True, align=PP_ALIGN.CENTER)
    rich(tf2, [(lab, False, INK)], size=11, align=PP_ALIGN.CENTER)
    x += 3.23

tb, tf = textbox(s, LM, 3.24, CW, 0.30)
rich(tf, [("제외 사유  ", True, GREY),
          ("무늬가 없음 218  ·  그런 옷이 아님 210  ·  색이 다름 10", False, INK)],
     size=11.5, first=True)

tb, tf = textbox(s, LM, 3.64, CW, 0.40)
sec_head(tf, "여기서 부딪힌 문제 — 거르는 순서", first=True)

TRACK_X, TRACK_W = 2.30, 7.20
for lab, pct, col, y_, note in [
        ("이전", 33.4, LGREY, 4.28, "다 만들고 나서 사진 없는 걸 버림"),
        ("현재", 97.3, NAVY,  5.10, "만들 때부터 사진 있는 것만 후보로")]:
    tb, tf = textbox(s, LM, y_ - 0.28, 7.0, 0.26)
    rich(tf, [(lab, True, RED if col == LGREY else NAVY),
              ("   " + note, False, GREY)], size=12, first=True)
    card(s, TRACK_X, y_, TRACK_W, 0.42, fill=RGBColor(0xF2, 0xF2, 0xF2),
         line=None, radius=0.06)
    fw = TRACK_W * pct / 100
    card(s, TRACK_X, y_, fw, 0.42, fill=col, line=None, radius=0.06)
    if pct > 60:
        tb, tf = textbox(s, TRACK_X + fw - 1.42, y_ + 0.05, 1.28, 0.32)
        rich(tf, [(f"{pct}%", True, WHITE)], size=16, first=True, align=PP_ALIGN.RIGHT)
    else:
        tb, tf = textbox(s, TRACK_X + fw + 0.16, y_ + 0.05, 1.4, 0.32)
        rich(tf, [(f"{pct}%", True, GREY)], size=16, first=True)
tb, tf = textbox(s, 7.10, 4.76, 2.40, 0.32)
rich(tf, [("약 3배", True, RED)], size=16, first=True, align=PP_ALIGN.RIGHT)
tb, tf = textbox(s, LM, 5.62, CW, 0.26)
rich(tf, [("사진 네 장이 모두 있는 문항의 비율", False, GREY)], size=11, first=True)

card(s, LM, 6.02, CW, 1.02, fill=PALER, line=LIGHT, radius=0.10)
tb, tf = textbox(s, LM + 0.20, 6.14, CW - 0.40, 0.80)
rich(tf, [("배운 것  ", True, NAVY),
          ("‘나중에 거르면 된다’가 사실이 아니었습니다. 버려진 문항이 한쪽에 몰려서 "
           "맞춰 놓은 균형이 깨졌습니다.", False, INK)], size=12.5, first=True, line=1.2)
rich(tf, [("그래서 ④의 결과를 ①의 입력으로 되돌렸습니다 — 조건을 새로 만든 게 "
           "아니라 앞으로 옮겼을 뿐입니다.", False, GREY)],
     size=11.5, space_before=4, line=1.2)
notes(s, "[3:45] 네 번째 단계입니다. 기계가 고른 사진을 사람이 한 칸씩 눈으로 "
         "확인했습니다. 1,661칸 중 1,223칸이 쓸 만했습니다. "
         "여기서 두 번째 문제를 만났습니다. 처음에는 문항을 다 만든 뒤 사진 없는 걸 "
         "버렸는데, 33%밖에 안 남았고 더 큰 문제는 버려진 게 한쪽에 몰려 균형이 "
         "깨진 것이었습니다. 그래서 사람이 확인한 목록을 첫 단계의 입력으로 "
         "되돌렸더니 97%로 올라가면서 균형도 유지됐습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 13 — ⑤ Eval : 모델 8종
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(13)
title(s, "⑤  Eval — 결과", sub="AI 8종 · 문항 2,550개 · 조건 5가지")

tb, tf = textbox(s, LM, 1.74, CW, 0.42)
rich(tf, [("따로 물어보면 잘합니다. ", True, INK),
          ("그런데 같이 물어보면 상황 쪽이 무너집니다", True, RED)], size=17, first=True)

pic(s, "fig_models.png", 0.72, 2.26, w=9.40)

tb, tf = textbox(s, LM, 5.90, CW, 0.40)
sec_head(tf, "읽는 법", first=True)

reads = [("예외가 없다", "여섯 모델 전부 같은 방향으로 떨어짐", MID),
         ("좋은 모델도 못 넘는다", "가장 좋은 모델의 둘 다 맞춘 비율 73%", RED),
         ("설계가 맞았다", "한쪽만 주면 다른 쪽은 딱 찍기 수준", NAVY)]
x = LM
for name, desc, col in reads:
    c = card(s, x, 6.32, 3.09, 0.76,
             fill=RGBColor(0xFF, 0xF3, 0xF3) if col == RED else WHITE,
             line=RED if col == RED else LIGHT, radius=0.10)
    tf2 = c.text_frame
    rich(tf2, [(name, True, col)], size=12.5, first=True, align=PP_ALIGN.CENTER)
    rich(tf2, [(desc, False, INK)], size=11, align=PP_ALIGN.CENTER)
    x += 3.23
notes(s, "[4:05] 마지막 단계, 결과입니다. AI 8종을 돌렸고 그중 여섯이 완주했습니다. "
         "왼쪽을 봐 주십시오. 상황만 알려줬을 때와 취향까지 알려줬을 때 상황 점수가 "
         "여섯 모델 전부 예외 없이 떨어집니다. 오른쪽은 둘 다 맞춘 비율인데, "
         "가장 좋은 모델도 73%입니다. 모델을 키운다고 해결되지 않았습니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 14 — ⑤ Eval : 사진 vs 글자
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(14)
title(s, "⑤  Eval — 못 하는 이유는 무엇인가", sub="보는 문제인가, 판단 문제인가")

lead(s, [("같은 모델에게 ", False, INK), ("같은 문항 2,550개", True, BLUE),
         ("를 사진 대신 글자로 주면 어떻게 될까?", False, INK)])

pic(s, "fig_modality.png", 0.96, 2.34, w=8.92)

card(s, LM, 4.90, CW, 0.74, fill=PALE, line=None, radius=0.10)
tb, tf = textbox(s, LM + 0.20, 5.04, CW - 0.40, 0.50)
rich(tf, [("글자로 줘도 83.8%에서 멈춥니다", True, NAVY),
          ("  —  남은 16%는 ", False, INK),
          ("보는 문제가 아니라 판단의 문제", True, RED), ("입니다.", False, INK)],
     size=14, first=True)

splits = [("10.6%p", "사진을 읽느라 잃는 몫", "옷을 보고 색·무늬를 알아보는 능력", MID),
          ("16.2%p", "글자로 줘도 남는 몫", "두 조건을 동시에 지키는 판단 능력", RED)]
x = LM
for v, name, desc, col in splits:
    c = card(s, x, 5.86, 4.66, 1.06,
             fill=WHITE if col == MID else RGBColor(0xFF, 0xF3, 0xF3),
             line=LIGHT if col == MID else RED, radius=0.10)
    tf2 = c.text_frame
    rich(tf2, [(v, True, col), ("   " + name, True, INK)], size=14,
         first=True, align=PP_ALIGN.CENTER)
    rich(tf2, [(desc, False, GREY)], size=11.5, align=PP_ALIGN.CENTER)
    x += 4.89
notes(s, "[4:20] 그러면 못 하는 이유가 사진을 못 봐서인지 판단을 못 해서인지 "
         "궁금했습니다. 같은 모델에게 같은 문항을 글자로 줘 봤습니다. "
         "73.1%에서 83.8%로 올랐습니다. 10.6%포인트는 사진을 읽느라 잃는 몫이고, "
         "나머지 16%는 글자로 줘도 남습니다. "
         "즉 병목이 보는 데만 있는 게 아니라 판단 자체에 있다는 뜻입니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 15 — ⑤ Eval : 어디가 특히 약한가
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(15)
title(s, "⑤  Eval — 어디가 특히 약한가", sub="Qwen3-VL 30B · 둘 다 알려준 조건")

lead(s, [("문항을 세 가지 기준으로 갈라 보면 ", False, INK),
         ("약한 곳이 분명하게 드러납니다", True, BLUE), (".", False, INK)])

pic(s, "fig_weakspots.png", 0.82, 2.32, w=9.20)

tb, tf = textbox(s, LM, 5.28, CW, 0.40)
sec_head(tf, "이 세 가지가 다음 연구의 출발점입니다", first=True)

weak = [("무늬가 가장 어렵다", "색은 80점인데 무늬는 63점 — 17점 차이", RED),
        ("격식 상황이 더 어렵다", "춥다·덥다보다 사회적 규범 쪽이 5점 낮음", MID),
        ("돌려 말하면 더 어렵다", "‘한겨울’처럼 간접적으로 주면 6점 하락", MID)]
y = 5.66
for name, desc, col in weak:
    card(s, LM, y, CW, 0.42, fill=RGBColor(0xFF, 0xF3, 0xF3) if col == RED else PALER,
         line=None, radius=0.10)
    tb, tf = textbox(s, LM + 0.22, y + 0.09, CW - 0.44, 0.26)
    rich(tf, [(name, True, col), ("      " + desc, False, INK)], size=12.5, first=True)
    y += 0.44
notes(s, "[4:35] 어디가 약한지도 갈라 봤습니다. 취향이 무늬일 때가 가장 어렵습니다. "
         "색은 80점인데 무늬는 63점입니다. 격식 같은 사회적 상황이 춥다·덥다보다 "
         "어렵고, 상황을 돌려 말하면 더 떨어집니다. "
         "이 세 가지가 다음 연구의 출발점입니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 15 — 기여
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(16)
title(s, "기여", sub="이번 인턴십에서 남긴 것")

contrib = [("문제를 나눌 수 있는 형태로 바꿈",
            "‘틀렸다’를 ‘어느 쪽에서 틀렸다’로 — 답 하나로 두 능력을 분해",
            "설계"),
           ("취향을 대화로 주는 조건을 새로 만듦",
            "값을 사진에만 담은 대화 48개 · 검증 48/48 · 돌연변이 6/6 검출",
            "신규"),
           ("가설을 8종 모델로 확인",
            "여섯 모델 전부에서 같은 방향 — 모델을 키워도 사라지지 않음",
            "실험"),
           ("병목의 위치를 갈라냄",
            "사진을 글자로 바꿔도 83.8% — 보는 문제가 아니라 판단의 문제",
            "분석"),
           ("누구나 다시 만들 수 있게 남김",
            "씨앗 고정 · 사람 검수 1,661칸 · 검증기를 생성기와 분리",
            "재현")]
y = 2.00
for name, desc, tag in contrib:
    card(s, LM, y, CW, 0.74, fill=PALER, line=LIGHT, radius=0.10)
    chip(s, LM + 0.20, y + 0.20, 0.62, 0.34, tag, MID, size=11)
    tb, tf = textbox(s, LM + 1.00, y + 0.13, CW - 1.24, 0.30)
    rich(tf, [(name, True, BLUE)], size=13.5, first=True)
    tb, tf = textbox(s, LM + 1.00, y + 0.44, CW - 1.24, 0.26)
    rich(tf, [(desc, False, INK)], size=11.5, first=True)
    y += 0.84

card(s, LM, 6.30, CW, 0.62, fill=WHITE, line=RED, radius=0.10)
tb, tf = textbox(s, LM + 0.20, 6.42, CW - 0.40, 0.40)
rich(tf, [("한계  ", True, RED),
          ("대화 조건은 만들고 검증까지 마쳤지만 ", False, INK),
          ("아직 모델에 태워 보지 못했습니다", True, INK),
          (" — 다음 단계입니다.", False, INK)], size=12.5, first=True)
notes(s, "[4:40] 기여는 다섯 가지입니다. 문제를 나눌 수 있는 형태로 바꿨고, "
         "취향을 대화로 주는 조건을 새로 만들었고, 가설을 여덟 모델로 확인했고, "
         "병목이 보는 데인지 판단인지 갈라냈고, 누구나 다시 만들 수 있게 남겼습니다. "
         "다만 대화 조건은 만들고 검증까지 마쳤지 아직 모델에 태워 보지는 "
         "못했습니다. 이게 바로 다음 단계입니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 16 — 향후 계획 · 소감
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide(17)
title(s, "향후 계획 · 소감")

tb, tf = textbox(s, LM, 1.80, 4.62, 0.4)
sec_head(tf, "남은 일", first=True)

plans = [("대화 조건 평가 실행", "만든 대화 48개를 모델에 태우기"),
         ("사진 라벨 검증 마무리", "관습 기반 35개 상황"),
         ("상용 모델 확대", "GPT · Gemini 계열 추가")]
y = 2.24
for name, desc in plans:
    card(s, LM, y, 4.62, 0.66, fill=PALER, line=LIGHT, radius=0.10)
    tb, tf = textbox(s, LM + 0.18, y + 0.12, 4.34, 0.44)
    rich(tf, [(name, True, BLUE)], size=13, first=True)
    rich(tf, [(desc, False, GREY)], size=11, line=1.1)
    y += 0.76

tb, tf = textbox(s, 5.57, 1.80, 4.62, 0.4)
sec_head(tf, "학업 · 진로 계획", first=True)

career = [("단기", "이 연구를 논문 투고 수준으로 정리"),
          ("중기", "AI 평가 방법론을 주제로 대학원 진학 준비"),
          ("장기", "‘모델이 무엇을 못하는지’를 재는 연구자")]
y = 2.28
for name, desc in career:
    chip(s, 5.57, y, 0.66, 0.30, name, MID, size=11)
    tb, tf = textbox(s, 6.38, y + 0.02, 3.81, 0.54)
    rich(tf, [(desc, False, INK)], size=12.5, first=True, line=1.2)
    y += 0.66

card(s, 5.57, 4.30, 4.62, 2.14, fill=PALE, line=None, radius=0.08)
tb, tf = textbox(s, 5.77, 4.44, 4.24, 1.88)
rich(tf, [("소감", True, NAVY)], size=13.5, first=True)
rich(tf, [("가장 크게 배운 것은 ", False, INK),
          ("실패한 결과를 그대로 남기는 법", True, INK), ("이었습니다.", False, INK)],
     size=12.5, space_before=7, line=1.25)
rich(tf, [("숫자가 나쁘면 숫자를 고치고 싶어지는데, 두 번 다 원인을 찾아 "
           "다시 설계해야 했습니다.", False, INK)], size=12.5, space_before=5, line=1.25)
rich(tf, [("성능보다 ", False, INK), ("측정이 옳은지", True, NAVY),
          ("를 먼저 의심하게 됐습니다.", False, INK)], size=12.5, space_before=5, line=1.25)

card(s, LM, 4.66, 4.62, 1.78, fill=WHITE, line=NAVY, radius=0.08)
tb, tf = textbox(s, LM + 0.18, 4.80, 4.30, 1.52)
rich(tf, [("전공 수업이 어디에 쓰였나", True, NAVY)], size=13, first=True)
for a, b in [("실험설계 · 통계", "두 축을 분리하는 설계"),
             ("컴퓨터비전", "검색 모델이 무엇을 놓치는지"),
             ("소프트웨어공학", "다시 만들 수 있는 파이프라인")]:
    rich(tf, [("· ", False, MID), (a, True, INK), ("   " + b, False, GREY)],
         size=11.5, space_before=6, line=1.15)
notes(s, "[4:50] 앞으로는 만든 대화를 실제로 모델에 태우고, 남은 사진 라벨 검증과 "
         "상용 모델 평가를 마무리해 논문 투고 수준으로 정리하려 합니다. "
         "소감을 말씀드리면, 가장 크게 배운 것은 실패한 결과를 그대로 남기는 "
         "법이었습니다. 숫자가 나쁘면 숫자를 고치고 싶어지는데, 두 번 다 원인을 "
         "찾아 다시 설계해야 했습니다. 성능보다 측정이 옳은지를 먼저 의심하게 "
         "됐고, 그게 전공 수업에서 배운 실험 설계를 제 손으로 써 본 경험이었습니다. "
         "감사합니다.")

# ═════════════════════════════════════════════════════════════════════════════
# 17 — Q&A
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
tb, tf = textbox(s, 0, 3.22, SW, 0.9)
rich(tf, [("Q & A", True, BLUE)], size=52, first=True, align=PP_ALIGN.CENTER)
tb, tf = textbox(s, 0, 4.22, SW, 0.4)
rich(tf, [("감사합니다.", False, GREY)], size=15, first=True, align=PP_ALIGN.CENTER)
tb, tf = textbox(s, 0, 4.66, SW, 0.34)
rich(tf, [("인공지능학과  2023480017  조현준   ·   VisAGI Lab", False, GREY)],
     size=12, first=True, align=PP_ALIGN.CENTER)
s.shapes.add_picture(LOGO, Inches(8.25), Inches(6.74), Inches(1.95), Inches(0.32))
notes(s, "[5:00] 감사합니다.")

prs.save(OUT)

# ── ship-time integrity gate ────────────────────────────────────────────────
import zipfile
from sanitize_template import verify

verify(OUT)
_z = zipfile.ZipFile(OUT)
_pres = _z.read("ppt/presentation.xml").decode("utf-8")
assert "<p:embeddedFontLst>" not in _pres, "embedded fonts leaked back into the deck"
assert not [n for n in _z.namelist() if n.startswith("ppt/fonts/")], "fntdata leaked back in"
_n_slides = len([n for n in _z.namelist() if n.startswith("ppt/slides/slide")])
_app = _z.read("docProps/app.xml").decode("utf-8")
import re as _re
_claim = _re.search(r"<Slides>(\d+)</Slides>", _app)
assert _claim is None or int(_claim.group(1)) == _n_slides, (
    f"docProps/app.xml claims {_claim.group(1)} slides, package has {_n_slides}")
_z.close()

print("saved:", OUT, "| slides:", _n_slides,
      f"| {os.path.getsize(OUT)/1e6:.2f} MB")
